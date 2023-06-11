import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Function to extract equations and graphs from a .docx file
def extract_math_content(file_path):
    document = Document(file_path)
    math_content = []

    for paragraph in document.paragraphs:
        if paragraph._element.xml.startswith('<w:pict'):
            math_content.append(paragraph._element.xml)

    return math_content

# Function to create a PowerPoint slide with math content
def create_ppt_slide(presentation, math_content):
    slide_layout = presentation.slide_layouts[1]  # Choose the desired slide layout

    slide = presentation.slides.add_slide(slide_layout)
    shapes = slide.shapes

    # Calculate the slide dimensions
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Create a text box to hold the math content
    left = Inches(1)
    top = Inches(1)
    width = slide_width - 2 * Inches(1)
    height = slide_height - 2 * Inches(1)

    textbox = shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # Add the math content to the text box
    p = text_frame.add_paragraph()
    for xml in math_content:
        p._element.xml = xml

# Main program
def main():
    # Specify the input .docx file path
    docx_file = "C:\\Users\Krishnendu Das\OneDrive\Desktop\DIGIPLUS\Sample(doc2ppt).docx"

    # Extract math content from the .docx file
    math_content = extract_math_content(docx_file)

    # Create a PowerPoint presentation
    presentation = Presentation()

    # Create a slide for each math content extracted
    for content in math_content:
        create_ppt_slide(presentation, [content])

    # Specify the output PowerPoint file path
    ppt_file = "output.pptx"

    # Save the PowerPoint presentation
    presentation.save(ppt_file)

    print(f"Math content extracted from {docx_file} and saved to {ppt_file}.")

# Run the program
if __name__ == "__main__":
    main()
